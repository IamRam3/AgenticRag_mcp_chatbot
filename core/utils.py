# core/utils.py

import os
import pandas as pd
import docx
import pptx
import fitz  # PyMuPDF

def parse_file(path: str, ext: str):
    if ext == ".pdf":
        return parse_pdf(path)
    elif ext == ".docx":
        return parse_docx(path)
    elif ext == ".csv":
        return parse_csv(path)
    elif ext in [".txt", ".md"]:
        return parse_txt(path)
    elif ext == ".pptx":
        return parse_pptx(path)
    return []

def parse_pdf(path):
    doc = fitz.open(path)
    return [page.get_text() for page in doc]

def parse_docx(path):
    document = docx.Document(path)
    return [para.text for para in document.paragraphs if para.text.strip()]

def parse_csv(path):
    df = pd.read_csv(path)
    return [row.to_string() for _, row in df.iterrows()]

def parse_txt(path):
    with open(path, "r", encoding="utf-8") as f:
        return f.read().splitlines()

def parse_pptx(path):
    prs = pptx.Presentation(path)
    slides = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
        slides.append(" ".join(texts))
    return slides
